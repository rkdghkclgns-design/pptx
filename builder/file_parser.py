"""Parse source files (.md, .txt, .pdf, .pptx) and extract text content."""

import os
from pathlib import Path

from PyPDF2 import PdfReader
from pptx import Presentation


def parse_text_file(path: str) -> str:
    """Read .md or .txt file as UTF-8 text."""
    return Path(path).read_text(encoding="utf-8")


def parse_pdf_file(path: str) -> str:
    """Extract text from all pages of a PDF."""
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def parse_pptx_file(path: str) -> str:
    """Extract text from all slides/shapes of a PPTX."""
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


PARSERS = {
    ".md": parse_text_file,
    ".txt": parse_text_file,
    ".pdf": parse_pdf_file,
    ".pptx": parse_pptx_file,
}


def parse_all_files(directory: str) -> str:
    """Parse all supported files in directory and return combined content."""
    combined = []

    for filename in sorted(os.listdir(directory)):
        ext = Path(filename).suffix.lower()
        parser = PARSERS.get(ext)
        if parser is None:
            continue

        filepath = os.path.join(directory, filename)
        try:
            content = parser(filepath)
            if content.strip():
                combined.append(f"=== {filename} ===\n{content}")
        except Exception as e:
            combined.append(f"=== {filename} ===\n[파일 파싱 오류: {e}]")

    return "\n\n".join(combined)
