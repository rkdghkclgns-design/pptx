"""Build PPTX presentations from structured slide data using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from schemas.slide_schema import SlideData

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x16, 0x20, 0x3A)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
TEXT_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_MUTED = RGBColor(0xA1, 0xA1, 0xAA)
TEXT_DARK = RGBColor(0x27, 0x27, 0x2A)
DIVIDER = RGBColor(0x3F, 0x3F, 0x46)

# Fonts (NanumGothic on Linux, fallback to Arial)
FONT_PRIMARY = "NanumGothic"
FONT_FALLBACK = "Arial"

# Margins
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
CONTENT_WIDTH = Inches(11.7)


def _font(name: str = FONT_PRIMARY) -> str:
    return name


def _add_background(slide, color: RGBColor = BG_DARK):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=None):
    """Helper to add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = _font(font_name or FONT_PRIMARY)
    p.alignment = alignment
    return txBox


def _add_bullets(slide, left, top, width, height, bullets: list[str],
                 font_size=16, color=TEXT_WHITE):
    """Add a bulleted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = _font()
        p.space_after = Pt(8)
        p.level = 0

    return txBox


def build_cover_slide(prs: Presentation, slide: SlideData):
    """Build a cover/title slide."""
    layout = prs.slide_layouts[6]  # Blank
    s = prs.slides.add_slide(layout)
    _add_background(s)

    # Accent line
    s.shapes.add_shape(
        1, MARGIN_LEFT, Inches(2.8), Inches(1.5), Pt(4)
    ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = ACCENT_CYAN

    # Title
    _add_text_box(s, MARGIN_LEFT, Inches(3.0), CONTENT_WIDTH, Inches(1.2),
                  slide.title, font_size=40, bold=True, color=TEXT_WHITE)

    # Subtitle
    if slide.subtitle:
        _add_text_box(s, MARGIN_LEFT, Inches(4.4), CONTENT_WIDTH, Inches(0.8),
                      slide.subtitle, font_size=20, color=TEXT_MUTED)

    # Speaker notes
    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_content_slide(prs: Presentation, slide: SlideData):
    """Build a standard content slide with title and bullets."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    # Title
    _add_text_box(s, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7),
                  slide.title, font_size=28, bold=True, color=TEXT_WHITE)

    # Divider line
    s.shapes.add_shape(
        1, MARGIN_LEFT, Inches(1.4), CONTENT_WIDTH, Pt(1)
    ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = DIVIDER

    y_offset = Inches(1.7)

    # Description
    if slide.description:
        _add_text_box(s, MARGIN_LEFT, y_offset, CONTENT_WIDTH, Inches(1.0),
                      slide.description, font_size=16, color=TEXT_MUTED)
        y_offset = Inches(2.8)

    # Bullets
    if slide.bullets:
        _add_bullets(s, MARGIN_LEFT, y_offset, CONTENT_WIDTH, Inches(4.0),
                     slide.bullets)

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_two_column_slide(prs: Presentation, slide: SlideData):
    """Build a two-column layout slide."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    _add_text_box(s, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7),
                  slide.title, font_size=28, bold=True, color=TEXT_WHITE)

    bullets = slide.bullets or []
    mid = len(bullets) // 2
    left_bullets = bullets[:mid] if mid > 0 else bullets
    right_bullets = bullets[mid:] if mid > 0 else []

    col_width = Inches(5.5)
    col_gap = Inches(0.7)

    # Left column card
    left_shape = s.shapes.add_shape(
        1, MARGIN_LEFT, Inches(1.6), col_width, Inches(5.0)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = BG_CARD
    left_shape.line.fill.background()

    if left_bullets:
        _add_bullets(s, Inches(1.2), Inches(2.0), Inches(4.8), Inches(4.2),
                     left_bullets, font_size=15)

    # Right column card
    right_left = MARGIN_LEFT + col_width + col_gap
    right_shape = s.shapes.add_shape(
        1, right_left, Inches(1.6), col_width, Inches(5.0)
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = BG_CARD
    right_shape.line.fill.background()

    if right_bullets:
        _add_bullets(s, right_left + Inches(0.4), Inches(2.0), Inches(4.8),
                     Inches(4.2), right_bullets, font_size=15)

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_three_cards_slide(prs: Presentation, slide: SlideData):
    """Build a three-card layout slide."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    _add_text_box(s, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7),
                  slide.title, font_size=28, bold=True, color=TEXT_WHITE)

    bullets = slide.bullets or []
    card_width = Inches(3.6)
    card_gap = Inches(0.45)

    for i in range(3):
        left = MARGIN_LEFT + (card_width + card_gap) * i
        shape = s.shapes.add_shape(1, left, Inches(1.6), card_width, Inches(5.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.fill.background()

        # Card content
        if i < len(bullets):
            _add_text_box(s, left + Inches(0.3), Inches(2.0),
                          card_width - Inches(0.6), Inches(4.2),
                          bullets[i], font_size=14, color=TEXT_WHITE)

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_table_slide(prs: Presentation, slide: SlideData):
    """Build a slide with a data table."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    _add_text_box(s, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7),
                  slide.title, font_size=28, bold=True, color=TEXT_WHITE)

    headers = slide.tableHeaders or []
    rows_data = slide.tableRows or []

    if not headers:
        return

    num_rows = len(rows_data) + 1
    num_cols = len(headers)

    table_shape = s.shapes.add_table(
        num_rows, num_cols,
        MARGIN_LEFT, Inches(1.6), CONTENT_WIDTH, Inches(5.0)
    )
    table = table_shape.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.font.name = _font()
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE

    # Data rows
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            if j >= num_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_WHITE
                p.font.name = _font()
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else BG_DARK

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_quote_slide(prs: Presentation, slide: SlideData):
    """Build a centered quote slide."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    quote_text = slide.description or slide.title
    _add_text_box(s, Inches(2.0), Inches(2.5), Inches(9.3), Inches(2.5),
                  f'"{quote_text}"', font_size=28, color=ACCENT_CYAN,
                  alignment=PP_ALIGN.CENTER)

    if slide.subtitle:
        _add_text_box(s, Inches(2.0), Inches(5.2), Inches(9.3), Inches(0.6),
                      f"- {slide.subtitle}", font_size=16, color=TEXT_MUTED,
                      alignment=PP_ALIGN.CENTER)

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_section_slide(prs: Presentation, slide: SlideData):
    """Build a section divider slide."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    # Section number/accent
    s.shapes.add_shape(
        1, Inches(5.5), Inches(2.5), Inches(2.3), Pt(4)
    ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = ACCENT_CYAN

    _add_text_box(s, Inches(2.0), Inches(3.0), Inches(9.3), Inches(1.2),
                  slide.title, font_size=36, bold=True, color=TEXT_WHITE,
                  alignment=PP_ALIGN.CENTER)

    if slide.subtitle:
        _add_text_box(s, Inches(2.0), Inches(4.4), Inches(9.3), Inches(0.8),
                      slide.subtitle, font_size=18, color=TEXT_MUTED,
                      alignment=PP_ALIGN.CENTER)

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def build_closing_slide(prs: Presentation, slide: SlideData):
    """Build a closing/thank-you slide."""
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    _add_background(s)

    _add_text_box(s, Inches(2.0), Inches(2.8), Inches(9.3), Inches(1.2),
                  slide.title, font_size=40, bold=True, color=TEXT_WHITE,
                  alignment=PP_ALIGN.CENTER)

    if slide.subtitle:
        _add_text_box(s, Inches(2.0), Inches(4.2), Inches(9.3), Inches(0.6),
                      slide.subtitle, font_size=18, color=TEXT_MUTED,
                      alignment=PP_ALIGN.CENTER)

    # Accent line
    s.shapes.add_shape(
        1, Inches(5.5), Inches(5.2), Inches(2.3), Pt(4)
    ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = ACCENT_CYAN

    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


# Map slide types to builder functions
SLIDE_BUILDERS = {
    "cover": build_cover_slide,
    "content": build_content_slide,
    "twoColumn": build_two_column_slide,
    "threeCards": build_three_cards_slide,
    "table": build_table_slide,
    "quote": build_quote_slide,
    "section": build_section_slide,
    "closing": build_closing_slide,
}


def build_pptx(slides: list[SlideData], output_path: str) -> str:
    """Build a complete PPTX presentation from slide data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides:
        builder = SLIDE_BUILDERS.get(slide_data.type, build_content_slide)
        builder(prs, slide_data)

    prs.save(output_path)
    return output_path
